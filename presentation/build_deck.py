"""Builds the AutoWorks stakeholder demo deck (.pptx).

Embeds app screenshots and the sample invoice PDF (as a double-clickable object).
Run:  python3 build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
PDF = os.path.join(HERE, "AutoWorks-Sample-Invoice.pdf")

# --- palette ---
DARK   = RGBColor(0x0F, 0x17, 0x2A)
SLATE  = RGBColor(0x1E, 0x29, 0x3B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
LBLUE  = RGBColor(0x3B, 0x82, 0xF6)
BG     = RGBColor(0xF8, 0xFA, 0xFC)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
TEXT   = RGBColor(0x0F, 0x17, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
GREEN  = RGBColor(0x05, 0x96, 0x69)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Calibri"


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    solid(sp, color)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text,size,color,bold,space_after,bullet}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        run = p.add_run()
        prefix = "•  " if ln.get("bullet") else ""
        run.text = prefix + ln["text"]
        run.font.name = FONT
        run.font.size = Pt(ln.get("size", 18))
        run.font.bold = ln.get("bold", False)
        run.font.color.rgb = ln.get("color", TEXT)
    return tb


def logo(slide, x, y, size=Inches(0.42), dark_text=False):
    box = rect(slide, x, y, size, size, BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.28
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "A"
    r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.name = FONT
    tb = slide.shapes.add_textbox(x + size + Inches(0.12), y - Inches(0.04),
                                  Inches(3), size + Inches(0.1))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "AutoWorks"
    r.font.bold = True; r.font.size = Pt(15)
    r.font.color.rgb = WHITE if not dark_text else TEXT
    r.font.name = FONT


def content_slide(kicker, title):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, SW, SH, WHITE)
    # header band
    logo(slide, Inches(0.55), Inches(0.42), dark_text=True)
    # accent line under header
    rect(slide, Inches(0.55), Inches(1.65), Inches(12.23), Pt(2), BORDER)
    # kicker + title
    textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.4),
            [{"text": kicker.upper(), "size": 12, "color": BLUE, "bold": True}])
    textbox(slide, Inches(0.55), Inches(1.18), Inches(12), Inches(0.55),
            [{"text": title, "size": 26, "color": TEXT, "bold": True}])
    return slide


def framed_picture(slide, path, x, y, w):
    pic = slide.shapes.add_picture(path, x, y, width=w)
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1)
    pic.shadow.inherit = False
    return pic


def footer(slide, n):
    textbox(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
            [{"text": "AutoWorks — Workshop Management System", "size": 9, "color": MUTED}])
    textbox(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.3),
            [{"text": str(n), "size": 9, "color": MUTED, "align": PP_ALIGN.RIGHT}])


# =========================================================
# 1. TITLE
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.22), SH, BLUE)
logo(s, Inches(0.9), Inches(0.85))
textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2),
        [
            {"text": "AutoWorks", "size": 54, "color": WHITE, "bold": True, "space_after": 2},
            {"text": "Workshop Management System", "size": 28, "color": LBLUE, "bold": True, "space_after": 14},
            {"text": "An end-to-end platform to run an automotive workshop — customers, "
                     "vehicles, work orders, and invoicing in one place.",
             "size": 16, "color": RGBColor(0xCB, 0xD5, 0xE1)},
        ])
textbox(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.6),
        [{"text": "Stakeholder Demo  ·  Phase 1 MVP", "size": 14,
          "color": RGBColor(0x94, 0xA3, 0xB8), "bold": True}])

# =========================================================
# 2. EXECUTIVE SUMMARY
# =========================================================
s = content_slide("Overview", "What is AutoWorks?")
textbox(s, Inches(0.55), Inches(2.0), Inches(7.2), Inches(4.5),
        [
            {"text": "A modern web application that digitizes the day-to-day operations of an "
                     "automotive repair workshop.", "size": 17, "color": TEXT, "space_after": 16},
            {"text": "Replaces paper job cards and spreadsheets", "size": 15, "color": TEXT, "bullet": True, "space_after": 8},
            {"text": "Single source of truth for customers & vehicles", "size": 15, "color": TEXT, "bullet": True, "space_after": 8},
            {"text": "Tracks every job from intake to payment", "size": 15, "color": TEXT, "bullet": True, "space_after": 8},
            {"text": "Generates invoices automatically from completed work", "size": 15, "color": TEXT, "bullet": True, "space_after": 8},
            {"text": "Real-time revenue & operations dashboard", "size": 15, "color": TEXT, "bullet": True},
        ])
# stat cards (stacked vertically on the right)
cards = [("5", "Core modules"), ("100%", "Web-based"), ("JWT", "Secure access")]
for i, (val, lbl) in enumerate(cards):
    top = Inches(2.0 + i * 1.45)
    card = rect(s, Inches(8.4), top, Inches(4.0), Inches(1.2), BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.08
    card.line.color.rgb = BORDER; card.line.width = Pt(1)
    textbox(s, Inches(8.7), top + Inches(0.16), Inches(3.4), Inches(0.95),
            [{"text": val, "size": 26, "color": BLUE, "bold": True, "space_after": 0},
             {"text": lbl, "size": 13, "color": MUTED}])
footer(s, 2)

# =========================================================
# 3. TECH STACK / ARCHITECTURE
# =========================================================
s = content_slide("Under the Hood", "Architecture & Technology")
layers = [
    ("Frontend", "Angular 19 · Angular Material · responsive SPA", LBLUE),
    ("REST API", "ASP.NET Core · Clean Architecture · FluentValidation", BLUE),
    ("Security", "ASP.NET Identity · JWT auth · role-based access", RGBColor(0x6D,0x28,0xD9)),
    ("Data", "Entity Framework Core · PostgreSQL 16", GREEN),
]
ty = 2.05
for name, desc, col in layers:
    bar = rect(s, Inches(0.55), Inches(ty), Inches(0.18), Inches(0.95), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    bar.adjustments[0] = 0.5
    card = rect(s, Inches(0.85), Inches(ty), Inches(11.9), Inches(0.95), BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.06
    card.line.color.rgb = BORDER; card.line.width = Pt(1)
    textbox(s, Inches(1.15), Inches(ty + 0.13), Inches(11.3), Inches(0.75),
            [{"text": name, "size": 16, "color": TEXT, "bold": True, "space_after": 2},
             {"text": desc, "size": 13, "color": MUTED}])
    ty += 1.18
footer(s, 3)

# =========================================================
# Feature slides (screenshot + bullets)
# =========================================================
def feature_slide(kicker, title, img, bullets, n):
    s = content_slide(kicker, title)
    framed_picture(s, os.path.join(ASSETS, img), Inches(0.55), Inches(1.95), Inches(8.1))
    lines = []
    for i, b in enumerate(bullets):
        lines.append({"text": b, "size": 14, "color": TEXT, "bullet": True, "space_after": 10})
    textbox(s, Inches(8.95), Inches(2.05), Inches(3.85), Inches(4.6), lines)
    footer(s, n)
    return s

feature_slide("Access", "Secure Login", "login.png",
              ["Email + password authentication",
               "JWT tokens issued on sign-in",
               "Role-based access (Admin, Manager, Technician)",
               "Clean, branded entry point"], 4)

feature_slide("At a glance", "Operations Dashboard", "dashboard.png",
              ["Revenue today / this week / this month",
               "Open work orders in progress",
               "Jobs completed this month",
               "Outstanding (unpaid) invoice totals",
               "Live KPIs the moment you log in"], 5)

feature_slide("Module", "Customers", "customers.png",
              ["Searchable customer directory",
               "Add new customers in seconds",
               "Contact details & vehicle counts",
               "Foundation for jobs and billing"], 6)

feature_slide("Module", "Vehicles", "vehicles.png",
              ["Every vehicle linked to its owner",
               "Make, model, year, plate, mileage, VIN",
               "Quick search across the fleet",
               "Add vehicles via guided form"], 7)

feature_slide("Module", "Work Orders", "work-orders.png",
              ["Track jobs from intake to completion",
               "Color-coded status workflow",
               "Filter by status",
               "One-click status updates",
               "Line-item pricing per job"], 8)

feature_slide("Workflow", "Creating a Work Order", "work-order-dialog.png",
              ["Pick customer, then their vehicle",
               "Add unlimited service line items",
               "Live estimated total as you type",
               "Capture customer notes",
               "Fast, validated data entry"], 9)

feature_slide("Module", "Invoices", "invoices.png",
              ["Generate invoices from completed jobs",
               "Automatic tax calculation",
               "Status lifecycle: Draft → Sent → Paid",
               "Track due dates & outstanding balances"], 10)

# =========================================================
# 11. INVOICE DOCUMENT  (the invoice IS the embedded PDF)
# =========================================================
s = content_slide("Deliverable", "Invoice Document")
inv_img = os.path.join(ASSETS, "invoice.png")
# The invoice graphic itself is the embedded PDF: frame == icon (square) so PowerPoint
# renders it 1:1 with no stretch. Double-clicking opens the full print-ready PDF.
embedded = False
try:
    s.shapes.add_ole_object(
        PDF, prog_id="AcroExch.Document.DC",
        left=Inches(0.55), top=Inches(1.95), width=Inches(4.95), height=Inches(4.95),
        icon_file=inv_img, icon_width=Inches(4.95), icon_height=Inches(4.95),
    )
    embedded = True
except Exception as e:  # pragma: no cover
    print("OLE embed failed, adding picture:", e)
    framed_picture(s, inv_img, Inches(0.55), Inches(1.95), Inches(4.95))

textbox(s, Inches(6.05), Inches(2.05), Inches(6.7), Inches(4.6),
        [
            {"text": "A print-ready invoice generated automatically from work-order data — "
                     "branded, itemized, with tax and totals.", "size": 16, "color": TEXT, "space_after": 16},
            {"text": "Customer & vehicle details", "size": 15, "color": TEXT, "bullet": True, "space_after": 10},
            {"text": "Itemized lines, subtotal, tax, total", "size": 15, "color": TEXT, "bullet": True, "space_after": 10},
            {"text": "Payment status & due dates", "size": 15, "color": TEXT, "bullet": True, "space_after": 10},
            {"text": "One-click generation from a completed job", "size": 15, "color": TEXT, "bullet": True, "space_after": 18},
            {"text": ("The real PDF is embedded in this slide — double-click the invoice to open it."
                      if embedded else "Print-ready PDF included alongside this deck."),
             "size": 13, "color": BLUE, "bold": True},
        ])
footer(s, 11)

# =========================================================
# 12. UI PROPOSAL & VISUAL THEME
# =========================================================
s = content_slide("Design", "UI Proposal & Visual Theme")


def col_header(slide, x, text):
    textbox(slide, Inches(x), Inches(2.0), Inches(3.9), Inches(0.35),
            [{"text": text.upper(), "size": 12, "color": MUTED, "bold": True}])


def swatch(slide, x, y, w, color, name, hexv, light=False):
    block = rect(slide, Inches(x), Inches(y), Inches(w), Inches(0.62), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    block.adjustments[0] = 0.18
    if light:
        block.line.color.rgb = BORDER
        block.line.width = Pt(1)
    textbox(slide, Inches(x), Inches(y + 0.66), Inches(w), Inches(0.5),
            [{"text": name, "size": 11, "color": TEXT, "bold": True, "space_after": 0},
             {"text": hexv, "size": 10, "color": MUTED}])


def chip(slide, x, y, w, h, text, bg, fg, line=None):
    sp = rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), bg, MSO_SHAPE.ROUNDED_RECTANGLE)
    sp.adjustments[0] = 0.5
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    tf = sp.text_frame
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = fg; r.font.name = FONT
    return sp


# --- Column 1: Color palette ---
col_header(s, 0.55, "Color Palette")
sw = [
    (BLUE, "Primary", "#2563EB", False),
    (RGBColor(0x1D, 0x4E, 0xD8), "Primary Dark", "#1D4ED8", False),
    (DARK, "Sidebar / Ink", "#0F172A", False),
    (GREEN, "Success", "#059669", False),
    (RGBColor(0xD9, 0x77, 0x06), "Warning", "#D97706", False),
    (BG, "Surface", "#F8FAFC", True),
]
sx = [0.55, 2.45]
sy = [2.45, 3.78, 5.11]
for i, (c, n, h, lt) in enumerate(sw):
    swatch(s, sx[i % 2], sy[i // 2], 1.75, c, n, h, lt)

# --- Column 2: Typography ---
col_header(s, 4.7, "Typography")
textbox(s, Inches(4.7), Inches(2.45), Inches(3.8), Inches(4.2),
        [
            {"text": "Inter", "size": 40, "color": TEXT, "bold": True, "space_after": 2},
            {"text": "Geometric, highly legible UI typeface", "size": 12, "color": MUTED, "space_after": 16},
            {"text": "Heading — 600", "size": 22, "color": TEXT, "bold": True, "space_after": 6},
            {"text": "Subheading — 500", "size": 16, "color": SLATE, "bold": True, "space_after": 6},
            {"text": "Body copy is set at a comfortable 14px with relaxed line-height for scanning dense data.",
             "size": 13, "color": TEXT, "space_after": 6},
            {"text": "CAPTION / META · 11px", "size": 10, "color": MUTED, "bold": True},
        ])

# --- Column 3: Components & patterns ---
col_header(s, 8.85, "Components")
# status badges
chip(s, 8.85, 2.45, 1.15, 0.42, "Paid", RGBColor(0xD1, 0xFA, 0xE5), GREEN)
chip(s, 10.10, 2.45, 1.45, 0.42, "In Progress", RGBColor(0xDB, 0xEA, 0xFE), BLUE)
chip(s, 8.85, 3.00, 1.55, 0.42, "Waiting Parts", RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xB4, 0x53, 0x09))
chip(s, 10.50, 3.00, 1.05, 0.42, "Overdue", RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0xDC, 0x26, 0x26))
# buttons
btn = rect(s, Inches(8.85), Inches(3.75), Inches(1.6), Inches(0.5), BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
btn.adjustments[0] = 0.18
tf = btn.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "+ Primary"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
ob = rect(s, Inches(10.55), Inches(3.75), Inches(1.6), Inches(0.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
ob.adjustments[0] = 0.18; ob.line.color.rgb = BORDER; ob.line.width = Pt(1.25)
tf = ob.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Secondary"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = TEXT; r.font.name = FONT
# sample card
card = rect(s, Inches(8.85), Inches(4.55), Inches(3.3), Inches(1.1), BG, MSO_SHAPE.ROUNDED_RECTANGLE)
card.adjustments[0] = 0.1; card.line.color.rgb = BORDER; card.line.width = Pt(1)
textbox(s, Inches(9.1), Inches(4.7), Inches(2.9), Inches(0.9),
        [{"text": "$12,480", "size": 22, "color": BLUE, "bold": True, "space_after": 0},
         {"text": "KPI card · rounded 12px · soft shadow", "size": 11, "color": MUTED}])

# --- Bottom: design principles (centered, fitted to width) ---
principles = ["8-px spacing grid", "Card-based surfaces", "12px rounded corners",
              "Soft elevation shadows", "Fully responsive"]
gap = 0.18
widths = [0.30 + 0.098 * len(t) for t in principles]
total = sum(widths) + gap * (len(principles) - 1)
px = (13.333 - total) / 2
for w, p_txt in zip(widths, principles):
    chip(s, px, 6.45, w, 0.42, p_txt, WHITE, SLATE, line=BORDER)
    px += w + gap
footer(s, 12)

# =========================================================
# 13. ROADMAP
# =========================================================
s = content_slide("What's next", "Roadmap")
done = ["Customer & vehicle management", "Work-order workflow & statuses",
        "Invoice generation + tax", "Analytics dashboard", "JWT security & roles"]
next_ = ["PDF export & email delivery", "Inline editing of records",
         "Technician scheduling / calendar", "Parts & inventory tracking",
         "Customer self-service portal"]
# Shipped column
hdr = rect(s, Inches(0.55), Inches(2.0), Inches(5.95), Inches(0.55), GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
hdr.adjustments[0] = 0.25
textbox(s, Inches(0.8), Inches(2.06), Inches(5.5), Inches(0.45),
        [{"text": "SHIPPED  (Phase 1 MVP)", "size": 13, "color": WHITE, "bold": True}])
textbox(s, Inches(0.7), Inches(2.75), Inches(5.7), Inches(4),
        [{"text": t, "size": 14, "color": TEXT, "bullet": True, "space_after": 11} for t in done])
# Next column
hdr = rect(s, Inches(6.8), Inches(2.0), Inches(5.95), Inches(0.55), BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
hdr.adjustments[0] = 0.25
textbox(s, Inches(7.05), Inches(2.06), Inches(5.5), Inches(0.45),
        [{"text": "PLANNED  (Next sprints)", "size": 13, "color": WHITE, "bold": True}])
textbox(s, Inches(6.95), Inches(2.75), Inches(5.7), Inches(4),
        [{"text": t, "size": 14, "color": MUTED, "bullet": True, "space_after": 11} for t in next_])
footer(s, 13)

# =========================================================
# 14. CLOSING
# =========================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.22), SH, BLUE)
logo(s, Inches(0.9), Inches(0.85))
textbox(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(2),
        [
            {"text": "Thank you", "size": 48, "color": WHITE, "bold": True, "space_after": 10},
            {"text": "Questions & live walkthrough", "size": 22, "color": LBLUE, "bold": True, "space_after": 18},
            {"text": "Demo login   admin@workshop.local  /  Admin123!", "size": 15,
             "color": RGBColor(0xCB, 0xD5, 0xE1)},
        ])

out = os.path.join(HERE, "AutoWorks-Demo.pptx")
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst), "embedded_pdf:", embedded)
